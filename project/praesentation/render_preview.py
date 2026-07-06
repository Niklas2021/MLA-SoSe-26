#!/usr/bin/env python3
# Rekonstruiert jede Folie der fertigen .pptx als Bild (Bilder an echter Position,
# Text mit echten Groessen/Farben, grob umgebrochen) -> Layout-Kontrolle ohne PowerPoint.
import io
import os
import textwrap
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle
from PIL import Image

EMU = 914400.0
SW, SH = 13.333, 7.5
HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "_layout")
os.makedirs(OUT, exist_ok=True)
prs = Presentation(os.path.join(HERE, "cuTile_Auto-Tuner.pptx"))


def rgb(c):
    return "#%02x%02x%02x" % (c[0], c[1], c[2])


def bg_color(slide):
    try:
        return rgb(slide.background.fill.fore_color.rgb)
    except Exception:
        return "white"


def fill_color(sh):
    try:
        if sh.fill.type == 1:
            return rgb(sh.fill.fore_color.rgb)
    except Exception:
        pass
    return "none"


def line_color(sh):
    try:
        if sh.line.width and sh.line.width > 0:
            return rgb(sh.line.color.rgb)
    except Exception:
        pass
    return "none"


for i, slide in enumerate(prs.slides):
    fig = plt.figure(figsize=(SW, SH), dpi=110)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, SW)
    ax.set_ylim(0, SH)
    ax.invert_yaxis()
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), SW, SH, facecolor=bg_color(slide), zorder=0))
    for sh in slide.shapes:
        L, T, W, H = sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU
        st = int(sh.shape_type) if sh.shape_type is not None else -1
        if st == 13:      # PICTURE
            im = Image.open(io.BytesIO(sh.image.blob))
            ax.imshow(im, extent=(L, L + W, T + H, T), zorder=3, aspect="auto")
        elif st == 1:     # AUTO_SHAPE
            fc, lc = fill_color(sh), line_color(sh)
            ax.add_patch(Rectangle((L, T), W, H, facecolor=fc,
                         edgecolor=(lc if lc != "none" else "none"),
                         linewidth=1.0, zorder=2))
        if sh.has_text_frame and sh.text_frame.text.strip():
            y = T + 0.06
            for p in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in p.runs)
                if not txt.strip():
                    continue
                run = next((r for r in p.runs if r.text.strip()), p.runs[0])
                sz = run.font.size.pt if run.font.size else 14
                try:
                    col = rgb(run.font.color.rgb)
                except Exception:
                    col = "black"
                bold = bool(run.font.bold)
                ha = {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}.get(p.alignment, "left")
                x = {"center": L + W / 2, "right": L + W - 0.08, "left": L + 0.08}[ha]
                cpl = max(6, int(W / (sz / 72.0 * 0.52)))
                for line in (textwrap.wrap(txt, cpl) or [""]):
                    ax.text(x, y, line, fontsize=sz, color=col, ha=ha, va="top",
                            fontweight="bold" if bold else "normal", zorder=4)
                    y += sz / 72.0 * 1.32
    fig.savefig(os.path.join(OUT, f"slide_{i+1:02d}.png"), dpi=110)
    plt.close(fig)

# Kontaktbogen (alle Folien auf einem Blatt)
imgs = sorted(os.path.join(OUT, f) for f in os.listdir(OUT) if f.startswith("slide_"))
cols, rows = 4, 5
fig, axes = plt.subplots(rows, cols, figsize=(cols * 3.4, rows * 2.0))
for ax in axes.flat:
    ax.axis("off")
for k, path in enumerate(imgs):
    ax = axes.flat[k]
    ax.imshow(Image.open(path))
    ax.set_title(f"Folie {k+1}", fontsize=8)
fig.tight_layout()
fig.savefig(os.path.join(OUT, "_contact_sheet.png"), dpi=110)
plt.close(fig)
print(f"{len(imgs)} Folien-Vorschauen + Kontaktbogen in {OUT}")
