#!/usr/bin/env python3
# Baut die PowerPoint-Praesentation (16:9) aus den Figures in figures/.
# Striktes Layout: jede Box hat feste Koordinaten (Zoll), Bilder werden per
# Seitenverhaeltnis in ihre Box zentriert -> nichts wird verzerrt oder verschoben.
# Detailinfos stehen in den Speaker-Notes, nicht auf den Folien.
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "cuTile_Auto-Tuner.pptx")

# --- Farben (wie in den Figures) ---
INK    = RGBColor(0x0b, 0x0b, 0x0b)
INK2   = RGBColor(0x3f, 0x3e, 0x3a)
MUTED  = RGBColor(0x89, 0x87, 0x81)
BLUE   = RGBColor(0x2a, 0x78, 0xd6)
ORANGE = RGBColor(0xeb, 0x68, 0x34)
NAVY   = RGBColor(0x0e, 0x1b, 0x2e)
WHITE  = RGBColor(0xff, 0xff, 0xff)
LIGHT  = RGBColor(0xf4, 0xf7, 0xfb)
FONT = "Calibri"

SW, SH = 13.333, 7.5           # Foliengroesse in Zoll
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_placed = []                   # (slide_idx, name, L, T, W, H) fuer den Bounds-Check


def _reg(idx, name, L, T, W, H):
    _placed.append((idx, name, L, T, W, H))


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set(run, size, color, bold=False, italic=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def textbox(s, L, T, W, H, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            idx=None):
    # lines: list of (text, size, color, bold, italic) ODER list davon je Absatz
    tb = s.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        runs = para if isinstance(para, list) else [para]
        for (text, size, color, bold, italic) in runs:
            r = p.add_run()
            r.text = text
            _set(r, size, color, bold, italic)
    if idx is not None:
        _reg(idx, "text", L, T, W, H)
    return tb


def rect(s, L, T, W, H, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(L), Inches(T), Inches(W), Inches(H))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def imsize(path):
    with Image.open(path) as im:
        return im.size


def image(s, name, L, T, W, H, idx=None):
    # zentriert das Bild seitenverhaeltnis-korrekt in die Box (L,T,W,H)
    path = os.path.join(FIG, name + ".png")
    w, h = imsize(path)
    ar, boxar = w / h, W / H
    if ar > boxar:
        nw, nh = W, W / ar
    else:
        nh, nw = H, H * ar
    l, t = L + (W - nw) / 2, T + (H - nh) / 2
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(nw), Inches(nh))
    if idx is not None:
        _reg(idx, name, l, t, nw, nh)
    return l, t, nw, nh


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


PAGE = [0]


def header(s, kicker, title, idx):
    # Eyebrow + Titel + Akzentlinie, einheitlich oben links
    PAGE[0] += 1
    textbox(s, 0.7, 0.42, 12.0, 0.28,
            [[(kicker.upper(), 12.5, BLUE, True, False)]], idx=idx)
    textbox(s, 0.68, 0.74, 12.0, 0.9,
            [[(title, 27, INK, True, False)]], idx=idx)
    rect(s, 0.72, 1.52, 2.1, 0.055, fill=BLUE)
    # Footer
    textbox(s, 0.7, 7.06, 7, 0.3, [[("cuTile Auto-Tuner · Becker-Klöser · Elagina",
            10, MUTED, False, False)]])
    textbox(s, 11.6, 7.06, 1.1, 0.3, [[(f"{PAGE[0]:02d}", 10, MUTED, False, False)]],
            align=PP_ALIGN.RIGHT)


def takeaway(s, text, idx):
    # Akzentbalken links + Kernaussage unten
    rect(s, 0.72, 6.42, 0.09, 0.44, fill=ORANGE)
    textbox(s, 0.95, 6.36, 11.6, 0.55,
            [[(text, 15, INK2, False, False)]], anchor=MSO_ANCHOR.MIDDLE, idx=idx)


def bullets(s, L, T, W, H, items, size=17, idx=None):
    tb = s.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.1
        r = p.add_run(); r.text = "▸  "; _set(r, size, BLUE, True)
        # it ist entweder ein String oder ein (text, color, bold)-Tupel
        if isinstance(it, tuple):
            t, c, b = it
            r2 = p.add_run(); r2.text = t; _set(r2, size, c, b)
        else:
            r2 = p.add_run(); r2.text = it; _set(r2, size, INK2, False)
    if idx is not None:
        _reg(idx, "bullets", L, T, W, H)
    return tb


def stat(s, L, T, W, value, label, color=BLUE):
    # kleine Stat-Kachel: grosse Zahl + Label
    textbox(s, L, T, W, 0.7, [[(value, 34, color, True, False)]],
            align=PP_ALIGN.CENTER)
    textbox(s, L, T + 0.72, W, 0.5, [[(label, 12.5, INK2, False, False)]],
            align=PP_ALIGN.CENTER)


# ============================================================
#  FOLIEN
# ============================================================
idx = 0

# --- 1 Titel ---
s = slide(NAVY)
rect(s, 0, 0, 0.35, SH, fill=BLUE)
rect(s, 0.9, 4.62, 2.0, 0.07, fill=ORANGE)
textbox(s, 0.9, 1.7, 11.5, 0.5, [[("GPU / cuTile · Projekt B", 15, BLUE, True, False)]])
textbox(s, 0.86, 2.25, 11.6, 2.3, [
    [("Auto-Tuning für", 44, WHITE, True, False)],
    [("Tensor-Kontraktionen", 44, WHITE, True, False)],
])
textbox(s, 0.9, 4.8, 11.5, 0.8,
        [[("Aus Einsum-String + Shapes automatisch eine gute cuTile-Tiling-Config finden — "
           "gemessen, nicht geraten.", 16.5, RGBColor(0xc3, 0xd4, 0xea), False, False)]])
textbox(s, 0.9, 6.35, 11.5, 0.36,
        [[("Niklas Becker-Klöser · Daria Elagina", 15, WHITE, True, False)]])
textbox(s, 0.9, 6.78, 11.5, 0.4,
        [[("cmk,ckn→cmn   ·   acspx,bspy→abcyx   ·   NVIDIA GB10 (DGX Spark)",
           12, MUTED, False, False)]])
notes(s, "Projekt B: Auto-Tuner fuer Tensor-Kontraktionen. Ziel: das Hand-Tuning aus A05/A06 "
         "automatisieren. Eingabe Einsum+Shapes, Ausgabe eine gute Tiling-Config, per do_bench "
         "gemessen. Zwei Testfaelle: batched Matmul (A05) und Tensor-Ring (A06). Alles auf der GB10.")

# --- 2 Recap / Problem ---
idx += 1
s = slide()
header(s, "Recap · Problem & Ziel", "Von Hand getunt skaliert nicht", idx)
bullets(s, 0.72, 1.9, 6.6, 4.0, [
    "A05 & A06: L2-optimale Aufteilung von Hand hergeleitet.",
    "Für jede neue Kontraktion / Shape / GPU: neu nachdenken.",
    "Ziel: Eingabe = Einsum + Shapes  →  Ausgabe = gute Config.",
    "Kein allgemeiner Tensor-Compiler — zwei Familien, Configs getunt.",
], size=18, idx=idx)
# rechte Karte: die zwei Testfaelle
rect(s, 7.7, 1.95, 4.9, 3.7, fill=LIGHT, line=RGBColor(0xdd,0xe4,0xec), line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 8.0, 2.15, 4.3, 0.5, [[("Zwei Testfälle", 15, INK, True, False)]])
textbox(s, 8.0, 2.75, 4.4, 1.2, [
    [("A05 · batched Matmul", 13.5, BLUE, True, False)],
    [("cmk, ckn → cmn", 17, INK, False, False)],
])
textbox(s, 8.0, 4.15, 4.4, 1.2, [
    [("A06 · Tensor-Ring", 13.5, ORANGE, True, False)],
    [("acspx, bspy → abcyx", 17, INK, False, False)],
])
takeaway(s, "Die Performance steckt in der Config, nicht im Kernel-Code — genau das automatisieren wir.", idx)
notes(s, "In A05/A06 haben wir die Tiling-Config von Hand begruendet. Das funktioniert, skaliert "
         "aber nicht auf neue Kontraktionen/Shapes/GPUs. Kernbeobachtung aus der Vorlesung: die "
         "Performance steckt in der Konfiguration (Tile-Groessen, Ausfuehrungsreihenfolge), nicht im "
         "Kernel-Code. Scope: kein allgemeiner Compiler, sondern ein kleiner Such-/Messloop.")

# --- 3 Divider Teil 1 ---
idx += 1
s = slide(NAVY)
rect(s, 0, 0, 0.35, SH, fill=BLUE)
textbox(s, 0.9, 2.72, 11.5, 0.4, [[("TEIL 1", 18, BLUE, True, False)]])
textbox(s, 0.86, 3.2, 11.6, 1.4, [[("Umsetzung & Architektur", 40, WHITE, True, False)]])
rect(s, 0.92, 4.5, 2.0, 0.07, fill=ORANGE)
textbox(s, 0.9, 4.75, 11.5, 0.5, [[("Pipeline · Suchraum · Kernel · Pruning · A06-Transfer",
        15, RGBColor(0xc3,0xd4,0xea), False, False)]])
notes(s, "Teil 1 (Person 1): wie der Tuner gebaut ist.")

# --- 4 Tiling-Config (fig_tiling) ---
idx += 1
s = slide()
header(s, "Vorlesungs-Recap", "Was ist eine Tiling-Config?", idx)
image(s, "fig_tiling", 0.6, 1.7, 8.0, 4.55, idx=idx)
bullets(s, 8.75, 2.0, 3.95, 4.1, [
    "Prim-Tiles M/N/K_PRIM: die mma-Kachel eines CTA.",
    "L2-Gruppe m_l2×n_l2: zeitlich nahe Blöcke → L2-Reuse.",
    "Exec-Reihenfolge PAR | SEQ | PRIM (K nie PAR).",
    "Variante A = Swizzle, B = SEQ-Loops.",
], size=15.5, idx=idx)
takeaway(s, "L2-Reuse entsteht durch die zeitliche Block-Gruppe, nicht durch eine räumliche Kachel.", idx)
notes(s, "Eine Config = Tile-Groessen + Ausfuehrungsreihenfolge. Prim-Tiles sind die mma-Kachel, die "
         "ein CTA rechnet. Die L2-Gruppe m_l2 x n_l2 sind benachbarte Bloecke, die zeitlich nah laufen, "
         "sodass A-Zeile und B-Spalte im L2 resident bleiben (Gegenstueck zu group_size_m im Triton-"
         "Matmul-Tutorial). Zwei Varianten: A = m_l2/n_l2 als PAR (Swizzle ueber die Block-ID), "
         "B = als SEQ-Loops im CTA.")

# --- 5 Pipeline (fig_pipeline) ---
idx += 1
s = slide()
header(s, "Überblick", "Die Tuner-Pipeline", idx)
image(s, "fig_pipeline", 0.6, 1.75, 12.1, 4.35, idx=idx)
takeaway(s, "generate → enumerate → prune → rank offline (reines Python); nur die Top-7 werden auf der GPU gemessen.", idx)
notes(s, "Fester Ablauf: generate_config erzeugt die Basic-Config; enumerate_candidates spannt den "
         "Suchraum auf (486); prune filtert statisch (342); rank sortiert per Kostenmodell (v2); tune "
         "kompiliert die Top-k, prueft gegen torch.einsum und misst mit do_bench; die Beste wird "
         "gecacht (Key = Einsum + Shapes + GPU-Modell). enumerate/prune/rank sind reines Python ohne "
         "GPU, nur tune braucht die Karte.")

# --- 6 Suchraum / Knoepfe ---
idx += 1
s = slide()
header(s, "Schritt 1", "Der Suchraum: die Knöpfe", idx)
bullets(s, 0.72, 1.95, 6.5, 3.6, [
    ("M_PRIM, N_PRIM ∈ {64, 128, 256}", INK, False),
    ("K_PRIM ∈ {32, 64, 128}", INK, False),
    ("M_L2, N_L2 ∈ {2, 4, 8}", INK, False),
    ("Variante ∈ {A, B}", INK, False),
], size=19, idx=idx)
# Formel-Karte + Check
rect(s, 7.55, 1.95, 5.05, 1.75, fill=LIGHT, line=RGBColor(0xdd,0xe4,0xec),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 7.75, 2.15, 4.7, 1.4, [
    [("3 · 3 · 3 · 3 · 3 · 2  =  ", 20, INK, True, False), ("486", 30, BLUE, True, False)],
    [("Kandidaten  (nicht die „81\" aus dem Pitch)", 12.5, INK2, False, False)],
], anchor=MSO_ANCHOR.MIDDLE)
rect(s, 7.55, 4.0, 5.05, 1.5, fill=RGBColor(0xea,0xf5,0xec),
     line=RGBColor(0xbf,0xe0,0xc8), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 7.75, 4.2, 4.7, 1.2, [
    [("✓  Akzeptanztest", 14, RGBColor(0x0c,0x83,0x0c), True, False)],
    [("Die handoptimierte A05-Config ist im Set — sonst könnte der Tuner sie nie finden.",
      13.5, INK2, False, False)],
], anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "Bewusst klein & hardware-sinnvoll gehalten; krumme Shapes werden hochgepaddet (PaddingMode.ZERO).", idx)
notes(s, "Der Suchraum ist bewusst klein. 3^5 * 2 = 486 (nicht die 81 aus dem Pitch, die zaehlten nur "
         "Tile-Kombis ohne asymmetrisches m_l2!=n_l2 und ohne die 2. Variante). Wichtiger Akzeptanztest: "
         "die handoptimierte A05-Config (128/128/64, 8x8, A) muss im Set sein. Krumme Shapes gehen nicht "
         "direkt durch split_dim (exakte Teilbarkeit), also padden wir hoch und nullen den Ueberhang.")

# --- 7 Kernel-Code (code_variant_a) ---
idx += 1
s = slide()
header(s, "Schritt 2 · Kernel", "Ein generischer Kernel — kein String-Codegen", idx)
image(s, "code_variant_a", 0.6, 1.7, 8.5, 4.5, idx=idx)
bullets(s, 9.3, 2.05, 3.45, 4.2, [
    "Tile-Größen als ct.Constant → JIT spezialisiert pro Wert.",
    "Keine fragilen String-Templates / exec().",
    "Variante A: Swizzle über die Block-ID.",
    "Korrektheit gegen torch.einsum (allclose).",
], size=14.5, idx=idx)
takeaway(s, "Compile ~0.4 s pro Config, auf der GB10 verifiziert — der ganze Ansatz hängt an dieser Spezialisierung.", idx)
notes(s, "Kern-Entscheidung: keine Kernel-Strings per exec(), sondern EIN generischer cuTile-Kernel pro "
         "Variante, Tile-Groessen als ct.Constant[int]. Der JIT spezialisiert pro Konstanten-Kombination "
         "(wie Triton constexpr). Auf der GB10 verifiziert: kompiliert, korrekt, spezialisiert pro Wert, "
         "~0.4 s Compile. Variante A dekodiert m_l2/n_l2 ueber die Block-ID (Swizzle).")

# --- 8 Pruning (fig_funnel + code_prune) ---
idx += 1
s = slide()
header(s, "Schritt 3 · Pruning", "Statisch eingrenzen — und wie weit wirklich", idx)
image(s, "fig_funnel", 0.5, 1.75, 6.15, 3.4, idx=idx)
image(s, "code_prune", 6.95, 1.78, 5.9, 3.0, idx=idx)
rect(s, 6.95, 5.0, 5.9, 1.24, fill=LIGHT, line=RGBColor(0xdd, 0xe4, 0xec),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 7.15, 5.08, 5.55, 1.08, [
    [("Was Filter 2 & 3 rechnen  (Formeln → nächste Folie)", 12, INK, True, False)],
    [("SMEM  (m·k + k·n)·2 B·2 Stages ≤ 100 KB   → 126 raus", 10.5, INK2, False, False)],
    [("Akku  m_prim·n_prim ≤ ½·65536 (fp32-Reg)   → 18 raus", 10.5, INK2, False, False)],
], anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "Nur 486→342: SMEM hängt nur an Prim-Größen — m_l2/n_l2 & Variante muss man messen (25 MB L2 killt die L2-Regel).", idx)
notes(s, "Vier Filter vor dem Kompilieren: (1) MMA-Teilbarkeit, (2) SMEM-Budget (hart), (3) Akku-Register, "
         "(4) Padding-Verschwendung. Ehrlich: A05 nur 486->342 (126 SMEM + 18 Register). Grund: SMEM haengt "
         "nur an den Prim-Groessen, nicht an m_l2/n_l2 oder der Variante — die kann statisches Pruning gar "
         "nicht anfassen. Und die L2-Reuse-Regel aus der Vorlesung greift auf der GB10 nicht (Working-Set "
         "~256 KB gegen 25 MB L2). Also verschiebt sich die Entscheidung auf die Messung.")

# --- 9 Mathematik (Pruning + Kostenmodell) ---
idx += 1
s = slide()
header(s, "Schritt 3 · Details", "Was Pruning und Kostenmodell konkret rechnen", idx)
image(s, "fig_math", 0.4, 1.72, 12.5, 4.5, idx=idx)
takeaway(s, "Alles aus device_props ausgelesen (SMEM / Register / Bandbreite / Peak) — nur Tile-Form, k_prim und L2-Gruppe variiert der Tuner.", idx)
notes(s, "Detail zur Pruning-Folie. PRUNING rechnet vier Dinge vor dem Compile: (1) alle Prim-Groessen "
         "Vielfache von 16 (fp16-Tensor-Cores); (2) Shared Memory pro Block = (m_prim·k_prim + "
         "k_prim·n_prim) · 2 Byte (fp16) · 2 (Double-Buffering) muss ins nutzbare SMEM (101376 − 1024 ≈ "
         "100 KB) passen → 126 Configs raus; (3) der fp32-Akku braucht m_prim·n_prim Register (1 pro "
         "Element), unter der halben Registerdatei (½·65536 = 32768) → 18 raus; (4) gepaddetes Volumen ≤ "
         "8× Original. RANKING schaetzt den DRAM-Traffic: A wird pro Gruppen-Spalte geladen, B pro "
         "Gruppen-Zeile, C einmal → Traffic ≈ [ M·K·ceil(N/(n_l2·n_prim)) + K·N·ceil(M/(m_l2·m_prim)) + "
         "M·N ] · 2 Byte. Groessere L2-Gruppe = kleinere ceil-Terme = weniger Traffic (das ist der "
         "L2-Reuse im Modell). Zeit = Traffic / Bandbreite, BW = mem_clock · (Bus/8) ≈ 273 GB/s (GB10, "
         "LPDDR). Die Roofline nimmt max(t_mem, t_compute) mit t_compute = 2·M·N·K / (Peak · util), "
         "Peak = SMs · Takt · FMAs · 2 ≈ 119 TFLOPS. Alle Hardware-Werte kommen aus device_props — "
         "deshalb ist das Modell portabel (auf der 3070 andere Zahlen, gleiche Formeln).")

# --- 9b Ranking-Mathematik: v2 vs roofline (progressiver Reveal, Teil 1) ---
# python-pptx kann keine Animationen -> "spaeteres Erscheinen" via Overlay-Folien:
# vier Folien, jede zeigt eine Stufe mehr. Durchklicken = eine Folie, die aufbaut.
# Alle vier teilen sich EINE Seitenzahl (RANK_PAGE), damit die Nummerierung stimmt.
RANK_PAGE = [0]


def _mono(s, L, T, W, H, lines, size=11.5, color=INK, idx=None):
    tb = s.shapes.add_textbox(Inches(L), Inches(T), Inches(W), Inches(H))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.06
        r = p.add_run()
        r.text = ln
        _set(r, size, color, False, False, "Consolas")
    if idx is not None:
        _reg(idx, "mono", L, T, W, H)
    return tb


def _panel(s, L, T, W, H, accent, title, idx=None):
    rect(s, L, T, W, H, fill=LIGHT, line=RGBColor(0xdd, 0xe4, 0xec), line_w=1.25,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, L + 0.28, T + 0.16, W - 0.5, 0.4, [[(title, 15.5, accent, True, False)]], idx=idx)
    rect(s, L + 0.3, T + 0.58, 0.85, 0.045, fill=accent)


def _badge(s, L, T, W, text, accent):
    rect(s, L, T, W, 0.5, fill=WHITE, line=accent, line_w=1.4,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, L + 0.15, T + 0.02, W - 0.3, 0.46, [[(text, 12.5, accent, True, False)]],
            anchor=MSO_ANCHOR.MIDDLE)


def ranking_math(level):
    global idx
    idx += 1
    s = slide()
    if level > 1:
        PAGE[0] = RANK_PAGE[0] - 1
    header(s, "Modell · Mathematik", "Zwei Kostenmodelle, zwei Physiken", idx)
    if level == 1:
        RANK_PAGE[0] = PAGE[0]

    textbox(s, 0.72, 1.62, 11.9, 0.4,
            [[("Input: der geprunte Pool (A05: 342 Configs). Pro Kandidat: Prim-Tiles "
               "m/n/k_prim + L2-Gruppe m_l2·n_l2.", 13, INK2, False, False)]])

    LX, RX, PT, PW, PH = 0.72, 6.83, 2.14, 5.78, 4.0
    _panel(s, LX, PT, PW, PH, BLUE, "v2 — Bandbreite  +  Register-Filter")
    _panel(s, RX, PT, PW, PH, ORANGE, "roofline — max(Memory, Compute)")

    if level >= 2:   # linkes Panel fuellen
        textbox(s, LX + 0.28, PT + 0.66, PW - 0.5, 0.3,
                [[("① Register-Filter (der billige Ausschluss)", 12.5, INK, True, False)]])
        _mono(s, LX + 0.28, PT + 0.98, PW - 0.5, 0.3,
              ["m_prim·n_prim ≤ 0.4·Regs  →  256-breite raus"], size=11)
        textbox(s, LX + 0.28, PT + 1.34, PW - 0.5, 0.3,
                [[("② Rank nach DRAM-Traffic / Bandbreite", 12.5, INK, True, False)]])
        _mono(s, LX + 0.28, PT + 1.66, PW - 0.5, 0.95, [
            "Traffic = M·K · ceil(N / n_l2·n_prim)",
            "        + K·N · ceil(M / m_l2·m_prim) + M·N",
            "t = Traffic·2B / Bandbreite     (kein L2!)",
        ], size=10.5, color=INK)
        textbox(s, LX + 0.28, PT + 2.66, PW - 0.5, 0.6,
                [[("= die Physik einer bandbreiten-limitierten GPU: kleines L2, "
                   "z. B. RTX 3070 (4 MB).", 11.5, INK2, False, True)]])
        _badge(s, LX + 0.28, PT + 3.35, PW - 0.56,
               "Spearman +0.38  ·  97.8 % Top-7  →  im Tuner", BLUE)

    if level >= 3:   # rechtes Panel fuellen
        _mono(s, RX + 0.28, PT + 0.66, PW - 0.5, 0.3,
              ["t = max( t_mem , t_compute )"], size=12, color=INK)
        textbox(s, RX + 0.28, PT + 1.02, PW - 0.5, 0.3,
                [[("t_mem: passt die L2-Gruppe ins L2?", 12, INK, True, False)]])
        _mono(s, RX + 0.28, PT + 1.34, PW - 0.5, 0.95, [
            "  ja   → nur Kaltladen (M·K+K·N+M·N)·2B",
            "  nein → Worst-Case wie v2",
            "t_compute = 2·M·N·K / (Peak · util)",
        ], size=10.5, color=INK)
        textbox(s, RX + 0.28, PT + 2.34, PW - 0.5, 0.9,
                [[("schaltet das Regime SELBST um (L2 aus device_props): kleines L2 → "
                   "memory,  GB10 25 MB → compute.", 11.5, INK2, False, True)]])
        _badge(s, RX + 0.28, PT + 3.35, PW - 0.56,
               "Spearman +0.50 (beste!)  ·  nur 85.5 % Top-7", ORANGE)

    if level >= 4:
        takeaway(s, "Compute-Regime: t_compute hängt nicht von m_l2/n_l2 ab → viele Gleichstände; "
                 "das feine L2-Reuse-Signal, das die Top-7 entscheidet, sieht roofline nicht. "
                 "Beste Korrelation ≠ bester Vorfilter.", idx)

    notes(s, "Beide Modelle bekommen DENSELBEN geprunten Pool. v2 wirft zuerst die Register-Fresser raus "
             "(m_prim·n_prim > 0.4·Regs, die 256-breiten Tiles) und rankt den Rest nach worst-case-DRAM-"
             "Traffic/Bandbreite -- das nimmt KEIN L2 an, groessere L2-Gruppe = weniger Nachladen. Das ist "
             "exakt die Physik einer bandbreiten-limitierten GPU mit kleinem L2 (RTX 3070, 4 MB) -> auf der "
             "GB10 (25 MB) ist es 'falsche Physik', trotzdem ist es als VORFILTER top (Spearman +0.38, 97.8 % "
             "Top-7). roofline rechnet stattdessen max(t_mem_L2, t_compute): t_mem_L2 ist L2-bewusst (passt die "
             "Gruppe ins L2, nur Kaltladen, sonst worst-case), t_compute = 2MNK/(Peak·util). Es schaltet das "
             "Regime selbst per L2-Groesse um (portabel) und korreliert global am besten (+0.50) -- aber im "
             "Compute-Regime der GB10 haengt t_compute gar nicht von m_l2/n_l2 ab, viele Gleichstaende, und das "
             "L2-Reuse-Signal 2. Ordnung, das die Top-7 entscheidet, sieht der Compute-Term nicht -> nur 85.5 % "
             "Top-7. Merksatz: bessere Korrelation != besserer Vorfilter. (bw ganz ohne Register-Filter: +0.03.)")


ranking_math(1)
ranking_math(2)
ranking_math(3)
ranking_math(4)

# --- 10 A06-Erweiterung (code_ring_a) ---
idx += 1
s = slide()
header(s, "Erweiterung", "A06: eine zweite Struktur-Familie", idx)
image(s, "code_ring_a", 0.6, 1.7, 8.4, 4.5, idx=idx)
bullets(s, 9.2, 2.05, 3.55, 4.2, [
    "Nicht „doppeltes K\" — die Batch-Topologie ist das Problem.",
    "A05: geteilter Batch c. A06: unabhängige a,c (A) und b (B).",
    "→ zweiter Kernel-Typ (Ring), kein Umbau.",
    "Der Tuner sucht Configs, nicht Kernel — A05-Pfad bleibt bitgleich.",
], size=14.5, idx=idx)
takeaway(s, "Wie cuBLAS/CUTLASS (Template-Menge) oder Triton (autotune pro @jit): neue Topologie = neues Template, Configs getunt.", idx)
notes(s, "A06-Ring acspx,bspy->abcyx: zwei Reduktionen (s,p), mehrere Output-Dims. Der eigentliche Knackpunkt "
         "ist die Batch-Topologie: A05 hat einen geteilten Batch c, A06 unabhaengige Batches (a,c nur in A, "
         "b nur in B). Der A05-Kernel indiziert A und B mit demselben c_idx — A06 kann das nicht ausdruecken. "
         "Deshalb ein zweiter Kernel-Typ (Ring-Kernel) mit Per-Tile-permute und aeusserer SEQ-Schleife ueber s. "
         "Enumerator 243 -> Pruning 171 (nur Variante A). Der Single-M/N/K-Pfad (A05) bleibt bitgleich.")

# --- 10 Die zwei Templates (Scope) ---
idx += 1
s = slide()
header(s, "Umsetzung · Config-Aufbau", "Wie wir die Reihenfolge der Dimensionen definiert haben", idx)
image(s, "fig_exec_order", 0.6, 1.66, 12.1, 4.55, idx=idx)
takeaway(s, "Jede Config ist by construction gültig: die Reihenfolge PAR │ SEQ │ PRIM steht fest (verify), der Tuner variiert nur die Split-Größen und A/B.", idx)
notes(s, "So haben wir die Reihenfolge definiert (verify + make_executable): die Dimensionen stehen immer "
         "in der Ordnung PAR links, dann SEQ, dann PRIM rechts. Regeln aus verify(): (1) keine K-Dim darf "
         "PAR sein (Reduktion nicht parallel), (2) alle SEQ links von allen PRIM, (3) alle PAR links von "
         "allen SEQ, (4) die rechtesten Dims sind PRIM und muessen je >=1 M-, N- und K-Dim enthalten (das "
         "ist die mma-Kachel). Der Tuner splittet jede M/N-Dim in l2_outer x l2 x prim und K in outer x "
         "prim, markiert die letzte M/N/K-Dim als PRIM und legt den Rest ab: m_l2/n_l2 werden in Variante A "
         "als PAR (Swizzle ueber die Block-ID) gesetzt, in Variante B als SEQ-Loops im CTA. A06 legt "
         "zusaetzlich die unabhaengigen Batches (a,c,b) als PAR und die zweite Reduktion s als SEQ ab. "
         "Kurz: der Tuner sucht nur die Groessen und A/B — die Struktur/Reihenfolge ist fix, verify() "
         "garantiert die Gueltigkeit. Den konkreten Kernel-Code dazu haben wir auf Folie 7 (A05) und 9 (A06).")

# --- 11 Divider Teil 2 ---
idx += 1
s = slide(NAVY)
rect(s, 0, 0, 0.35, SH, fill=BLUE)
textbox(s, 0.9, 2.72, 11.5, 0.4, [[("TEIL 2", 18, BLUE, True, False)]])
textbox(s, 0.86, 3.2, 11.6, 1.4, [[("Evaluation & Ergebnisse", 40, WHITE, True, False)]])
rect(s, 0.92, 4.5, 2.0, 0.07, fill=ORANGE)
textbox(s, 0.9, 4.75, 11.5, 0.5, [[("Gegen Handkernel · gegen cuBLAS · über Shapes & GPUs",
        15, RGBColor(0xc3,0xd4,0xea), False, False)]])
notes(s, "Teil 2 (Person 2): wie gut das Ganze misst.")

# --- 11 Benchmark-Setup ---
idx += 1
s = slide()
header(s, "Evaluation", "Benchmark-Setup: acht Shape-Regime je Familie", idx)
image(s, "fig_regimes", 0.5, 1.78, 7.9, 4.3, idx=idx)
bullets(s, 8.55, 2.05, 4.2, 4.0, [
    "Pro Shape: Tuner vs. Default (8×8) vs. torch.einsum; A06 auch vs. Handkernel.",
    "fp16 rein, fp32 Akku; jede Config gegen torch.einsum geprüft (allclose).",
    "8×342 (A05) + 8×171 (A06): 0 Fehlschläge, inkl. Padding-Pfad (krumm).",
    "GPU: NVIDIA GB10 — 48 SMs, 25 MB L2, integrierter LPDDR.",
], size=14.5, idx=idx)
takeaway(s, "Acht Regime decken square / rechteckig / K-Extreme / unteilbar / Batch ab — fair gegen drei Referenzen, alles korrekt.", idx)
notes(s, "Gemessen wird pro Shape: getunter Kernel, Default (naive 8x8), torch.einsum (cuBLAS), bei A06 "
         "zusaetzlich der Handkernel (46.5). fp16 rein, fp32 Akku. Jede Config wird gegen torch.einsum "
         "geprueft (allclose rtol=1e-2/atol=1e-1). Ergebnis: 8x342 (A05) + 8x171 (A06), alle korrekt, "
         "0 Fehlschlaege, inkl. der unteilbaren krumm-Shapes. GB10: 48 SMs, 25 MB L2, integriert.")

# --- 12 A05-Ergebnisse ---
idx += 1
s = slide()
header(s, "Ergebnisse · A05", "Der Tuner bestätigt die Handarbeit", idx)
image(s, "fig_a05_bars", 0.5, 1.7, 9.2, 4.5, idx=idx)
bullets(s, 9.85, 2.1, 2.9, 4.15, [
    ("Tuner-Pick (top-7) ≈ Default → bestätigt die Hand.", INK2, False),
    ("Bench-Best (von 342) nur knapp drüber; krumm +58 %.", INK2, False),
    ("Aber: selbst die Bench-Best bleibt auf GEMM unter cuBLAS.", INK2, False),
], size=13.5, idx=idx)
takeaway(s, "Warum verliert selbst unsere Bench-Best gegen torch? cuBLAS ist eine gereifte GEMM-Library (hand-optimiertes SASS, Split-K, Pipelining) — das schlägt ein einfaches Template.", idx)
notes(s, "Vier Balken: Default (naive 8x8), Tuner-Pick (bester der Modell-Top-7 = was der Tuner praktisch "
         "liefert), Bench Best (bestes von 342 gemessenen Configs = Voll-Sweep-Optimum), torch (cuBLAS). "
         "Auf regulaeren GEMMs ist der Tuner-Pick ~= Default (bestaetigt das Handtuning); die Bench-Best "
         "liegt nur ~1-3 % drueber, nur bei krumm deutlich (26.6 -> 41.8, +58 %). ENTSCHEIDEND fuer die "
         "Frage 'schlaegt torch unsere Bench-Best?': ja, auf 7/8 GEMM-Shapes liegt torch ueber unserer "
         "absoluten Bench-Best (nur die hand-getunte a05-Referenz zieht knapp vorbei). Der Grund ist nicht "
         "die Config-Suche, sondern die Kernel-Reife: cuBLAS nutzt hand-optimiertes SASS-Assembly, Split-K, "
         "ausgefeiltes Pipelining/Scheduling — Reife, die unser generisches ct.Constant-Template nicht "
         "erreicht. Das ist ehrlich und erwartbar; cuBLAS auf GEMM zu schlagen war nie das Ziel.")

# --- 13 cuBLAS / torch ehrlich ---
idx += 1
s = slide()
header(s, "Ergebnisse · ehrliche Baseline", "Tuner vs. cuBLAS / torch.einsum", idx)
image(s, "fig_tuner_vs_torch", 0.7, 1.68, 7.2, 4.6, idx=idx)
bullets(s, 8.2, 2.05, 4.55, 4.3, [
    "GEMM (A05): cuBLAS gewinnt (Tuner ~77 % von torch) — erwartbar.",
    "Ring (A06): im Mittel ~1.17× vorn, sehr shape-abhängig.",
    "Großer Gewinn, wo torchs Pfad schlecht ist (tall 3.95×, large_k 2.65×).",
    "A06_TORCH=16.18 aus dem Assignment ist veraltet → frisch 60.2.",
], size=14.5, idx=idx)
takeaway(s, "Der Wert ist nicht „schneller als alles\", sondern: ohne Handarbeit über beliebige Kontraktionen brauchbar — und stark, wo die Library keinen guten Pfad hat.", idx)
notes(s, "Auf reinem GEMM gewinnt cuBLAS (Tuner ~77 % von torch, geom. Mittel) — erwartbar, cuBLAS ist eine "
         "gereifte Library. Auf Ring-Shapes ist der Tuner im Mittel leicht vorn (~1.17x geom.), sehr "
         "shape-abhaengig: gross wo torchs Pfad schlecht ist (tall 3.95x, large_k 2.65x), schwach wo gut "
         "(krumm 0.38x). Ehrliche Korrektur: der Assignment-Wert A06_TORCH_EINSUM=16.18 ist veraltet; frisch "
         "macht dieselbe Shape 60.22 TFLOPS, gleichauf mit dem Tuner.")

# --- 14 A06-Ergebnisse (bars + ladder) ---
idx += 1
s = slide()
header(s, "Ergebnisse · A06", "Transfer gelingt — Tuner schlägt den Handkernel", idx)
image(s, "fig_a06_bars", 0.4, 1.72, 8.0, 4.5, idx=idx)
image(s, "fig_a06_ladder", 8.45, 1.72, 4.4, 4.5, idx=idx)
takeaway(s, "Anders als GEMM: unsere Bench-Best schlägt torch, wo dessen Ring-Pfad schlecht ist (tall, large_k); der Tuner schlägt zudem den Handkernel (+29 %, aus k_prim=32).", idx)
notes(s, "Referenz-Shape: Tuner ~60 gegen Handkernel 46.5 (+29 %), gegen die mismatchte 8x8-Default 26.3 "
         "(2.29x — aber das ist die falsche Messlatte), gleichauf mit frischem torch 60.2. Der Gewinn kommt "
         "aus k_prim=32: der Handkernel nahm p=64 als einen mma-Tile, der Tuner teilt in zwei 32er-Kacheln. "
         "Der Handkernel-Balken ist das feste Referenz-Tiling (2x3) auf jede Shape gelegt -- es passt nur "
         "auf der Referenz gut, liegt bei 5/8 Shapes sogar unter dem 8x8-Default, waehrend der Tuner in allen "
         "8 Regimen davor liegt. Ueber alle 8 Ring-Shapes 1.10-2.47x ueber Default, alle 171 Configs je Shape korrekt inkl. Padding. "
         "Vierter Balken (Bench Best = bestes von 171): schlaegt torch dort, wo dessen Ring-Pfad schlecht "
         "ist (square, tall, wide, large_k), verliert wo torch einen guten Pfad findet (small_k, krumm, "
         "batch). Anders als bei A05 (wo cuBLAS immer gewinnt) ist die Ring-Familie also der Fall, in dem "
         "sich der eigene Kernel lohnt.")

# --- 15 Top-k Stufen ---
idx += 1
s = slide()
header(s, "Modell · Praxis", "Wie oft ist das Optimum in den Top-k?", idx)
image(s, "fig_topk_curve", 0.6, 1.7, 8.6, 4.55, idx=idx)
bullets(s, 9.4, 2.15, 3.35, 4.0, [
    ("top-7 · ~3 s · ≥ 95 %", INK, True),
    ("top-45 · ~18 s · ≥ 99 %", INK, True),
    ("Voll-Sweep · ~3 min · 100 %", INK, True),
    ("Exakt-Beste sitzt tief — Spitzenfeld liegt im Messrauschen.", INK2, False),
], size=15, idx=idx)
takeaway(s, "Als Vorauswahl reicht das Modell fast immer; als exakter Treffer selten — top-7 ist der Sweet Spot.", idx)
notes(s, "Ground Truth = Voll-Sweep. Frage: zieht das Modell die real beste Config in die Top-k? Drei Stufen: "
         "top-7 (~3 s) -> im Schnitt 97 % des Optimums; top-45 (~18 s) -> 99 %; Voll-Sweep (~3 min) -> 100 %. "
         "Die exakt Beste erwischt top-7 meist nicht (sie sitzt im Modell tief, das Spitzenfeld liegt "
         "innerhalb ~3 % im Messrauschen). Kostet aber praktisch nichts an Performance.")

# --- 16 Ranking-Modell (Payoff-Scatter) ---
idx += 1
s = slide()
header(s, "Modell", "Besserer Ranker ≠ besserer Vorfilter", idx)
image(s, "fig_ranking_models", 0.6, 1.7, 8.4, 4.55, idx=idx)
bullets(s, 9.2, 2.1, 3.55, 4.1, [
    ("bw pur: +0.03 hier — aber die richtige Physik für kleine-L2-GPUs (3070).", INK2, False),
    ("v2 (bw + Reg-Filter): +0.38 · 97.8 % top-7 → im Tuner verbaut.", INK2, False),
    ("roofline (L2-bewusst): +0.50 beste Korrelation, nur 85.5 % top-7.", INK2, False),
], size=14, idx=idx)
takeaway(s, "Höchste Korrelation ≠ beste Top-k-Ausbeute — für die Praxis zählt der Vorfilter, also v2.", idx)
notes(s, "Payoff der Mathe-Folie in einem Bild: x = Korrelation (Spearman), y = Top-7-Ausbeute. bw pur rankt "
         "schlecht (Spearman +0.03) — falsche Physik auf der GB10 (25 MB L2, fast alles resident), aber genau "
         "das Modell, das man auf einer bandbreitenlimitierten Karte mit kleinem L2 (3070, 4 MB) nehmen wuerde. "
         "v2 (bw + Register-Filter) +0.38 und 97.8 % top-7 -> im Tuner verbaut. roofline (max(memory,compute), "
         "L2-bewusst) korreliert am besten (+0.50), ist aber schlechterer Vorfilter (85.5 %). Ueberleitung: "
         "genau dieser L2-Umschalter wird auf der 3070 gleich real (naechste/vorige Cross-GPU-Folie).")

# --- 17 Cross-GPU ---
idx += 1
s = slide()
header(s, "Ergebnisse · Cross-GPU", "Derselbe Hebel, aber andere Config pro GPU", idx)
image(s, "fig_crossgpu_lever", 0.55, 1.72, 12.2, 4.45, idx=idx)
takeaway(s, "Tuning hilft auf beiden Karten (Ø 1.36× GB10, 1.88× 3070) und die optimale Config ist in 16/16 Shapes verschieden — genau deshalb: GPU-spezifisch tunen + cachen.", idx)
notes(s, "Absolute TFLOPS der GB10 und 3070 zu vergleichen ist nicht sinnvoll (andere Peak-Leistung, "
         "Bandbreite, L2: 25 MB vs 4 MB). Vergleichbar ist der Optimierungshebel: der Speedup Tuner/Default "
         "pro Shape. Ergebnis: der Hebel wirkt auf beiden Karten, auf der 3070 sogar staerker (Ø 1.88x vs "
         "1.36x) — der aus dem GB10-Handtuning stammende 8x8-Default passt auf der 3070 schlechter, also "
         "holt der Tuner dort mehr raus. Und die gemessen beste Config ist in 16/16 Shapes zwischen den "
         "Karten verschieden (GB10 mag 128/128, die 3070 oft 64/128 oder 256/64 mit anderem k_prim/L2). "
         "Das ist das Argument fuer GPU-spezifisches Autotuning und warum das GPU-Modell im Cache-Key steht. "
         "Portabilitaet ist by construction (alles aus device_props) — jetzt auch auf einer zweiten, "
         "bandbreitenlimitierten GPU bestaetigt.")

# --- 18 Cross-GPU: Config-Wahl & Kosten ---
idx += 1
s = slide()
header(s, "Ergebnisse · Cross-GPU", "Was der Tuner pro GPU wählt — und was es kostet", idx)
image(s, "fig_config_table", 0.5, 1.72, 8.35, 4.4, idx=idx)
bullets(s, 9.0, 2.05, 3.75, 4.15, [
    "Optimum in 16/16 Shapes verschieden — pro GPU eine andere Config.",
    "Autotuner (v2, top-7) trifft ⌀ 96 % (GB10) / 90 % (3070) des Optimums.",
    "Kosten: GB10 ~3 s/Shape; 3070 (WSL-Compile) ~3 s–2 min — Compile dominiert.",
    "Roofline schaltet sein Regime selbst um; praktisch bleibt v2 der beste Vorfilter.",
], size=14, idx=idx)
takeaway(s, "Pro GPU von Hand zu tunen ist nicht machbar — das nimmt der gecachte Tuner ab, ohne auf eine Karte zu overfitten.", idx)
notes(s, "Die gemessen beste Config ist in 16/16 Shapes zwischen GB10 und 3070 verschieden: die GB10 mag "
         "grosse 128x128-Tiles (25 MB L2 vertraegt sie), die 3070 kleineres k_prim=32 und oft asymmetrische/"
         "64-breite Tiles. Der Autotuner (v2-Modell, nur top-7 gemessen) trifft im Schnitt ~96 % (GB10) bzw. "
         "~90 % (3070) des Optimums — auf der 3070 etwas unzuverlaessiger (einzelne Shapes wie krumm/square_1b "
         "fallen auf ~60-75 %). "
         "KOSTEN: eine Shape zu tunen (top-7 kompilieren+messen) kostet auf der GB10 ~3 s. Auf der 3070 lief "
         "alles unter WSL, wo der gcc-Compile viel langsamer ist — je nach Shape ~3 s bis ~2 min (ein voller "
         "Sweep dauerte dort bis ~1.9 h vs. ~5 min auf der GB10). Es ist ein Einmalkosten-Posten, der gecacht "
         "wird, und der Nutzen ist auf der 3070 groesser (Oe 1.88x) — fuer wiederverwendete Shapes lohnt es "
         "sich, fuer Einmalrechnungen auf der 3070 grenzwertig. "
         "BEWERTUNGSMETHODE: das Roofline-Modell schaltet sein Regime selbst per device_props um "
         "(max(memory,compute)). Interessant: die GB10 ist die bandbreiten-limitierte Karte (273 GB/s LPDDR, "
         "3/16 Shapes memory-bound), die 3070 bleibt compute-bound (448 GB/s GDDR6, 0/16). Der praktische "
         "Tuner (autotune.py) nutzt aber fest v2, weil v2 auf beiden Karten der bessere top-7-Vorfilter ist "
         "(roofline korreliert global besser: Spearman +0.57 vs v2 +0.19 auf der 3070, ist aber schlechterer "
         "Vorfilter: 83 % vs 91 %).")

# --- 19 Praxis / Cache ---
idx += 1
s = slide()
header(s, "Einordnung", "Wann sich Tuning lohnt: Cache & Amortisierung", idx)
bullets(s, 0.72, 1.95, 7.4, 4.0, [
    "Eine Kontraktion läuft ~8 ms; Tuning ~3 s ist ein Einmalposten.",
    "Im Netz sind Layer-Dims fix → millionenfach dieselbe Shape.",
    "Config-Cache: Key = Einsum + Shapes + GPU-Modell.",
    "Genau so machen es cuBLAS, Triton (autotune), torch.compile.",
], size=16.5, idx=idx)
stat(s, 8.5, 2.2, 4.1, "~3 s → 0 s", "erste getunte Shape → jede weitere aus dem Cache")
rect(s, 8.5, 4.0, 4.1, 1.6, fill=LIGHT, line=RGBColor(0xdd,0xe4,0xec),
     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, 8.72, 4.18, 3.7, 1.3, [
    [("End-to-end bestätigt", 13.5, INK, True, False)],
    [("16 Shapes getunt & gecacht, 2. Aufruf Cache-Hit, Top-7 bei ~95–100 %.",
      12.5, INK2, False, False)],
], anchor=MSO_ANCHOR.MIDDLE)
takeaway(s, "Break-even in Sekunden bei fester Shape; für stark dynamische Shapes hilft Bucketing/Padding.", idx)
notes(s, "Tuning lohnt sich ueber Wiederverwendung. Eine Kontraktion laeuft ~8 ms, die 3 s Tuning "
         "amortisieren sich, sobald dieselbe Shape oft laeuft — im Training/Inferenz der Normalfall (feste "
         "Layer-Dims). Config-Cache (cache.py + autotune.py), Key inkl. GPU-Modell. End-to-end bestaetigt: "
         "16 Shapes getunt & gecacht, 2. Aufruf Cache-Hit, Top-7-Picks bei ~95-100 %. Nur bei stark "
         "dynamischen Shapes geht die Rechnung nicht auf -> Bucketing/Padding.")

# --- 19b Live-Demo (2 Stufen: grau + "?" -> beim naechsten Klick das Bild) ---
DEMO_PAGE = [0]
GREY_FILL = RGBColor(0xe9, 0xec, 0xef)
GREY_LINE = RGBColor(0xcf, 0xd6, 0xde)


def demo_slide(reveal):
    global idx
    idx += 1
    s = slide()
    if reveal:
        PAGE[0] = DEMO_PAGE[0] - 1
    header(s, "Live-Demo", "Der Tuner malt ein Bild — live auf dem DGX", idx)
    if not reveal:
        DEMO_PAGE[0] = PAGE[0]

    if reveal:
        image(s, "fig_demo_preview", 0.55, 1.72, 7.55, 4.35, idx=idx)
    else:
        rect(s, 0.55, 1.772, 7.55, 4.246, fill=GREY_FILL, line=GREY_LINE, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        textbox(s, 0.55, 1.772, 7.55, 4.246, [[("?", 150, MUTED, True, False)]],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, idx=idx)
        textbox(s, 0.55, 5.35, 7.55, 0.5,
                [[("Tuner vs. Default — welches Bild kommt raus?", 15, MUTED, False, True)]],
                align=PP_ALIGN.CENTER)

    bullets(s, 8.3, 2.0, 4.55, 4.2, [
        "Full-HD-Plasma (2 MPixel) = ein großes GEMM: Y · X.",
        "8×8-Default vs. Tuner-Config — gleiche Mathe.",
        "Gleiches Bild, beide gegen torch geprüft.",
        "Krumme Shape → Default padded schlecht → Tuner gewinnt.",
    ], size=14, idx=idx)
    rect(s, 8.35, 5.55, 4.4, 0.6, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, 8.5, 5.62, 4.1, 0.46, [[("▶  python demo_paint.py", 15, WHITE, True, False)]],
            anchor=MSO_ANCHOR.MIDDLE)
    takeaway(s, "enumerate 486 → prune → Top-7 messen → schnellere Config, live. "
             "Gleiches Bild, aber der Tuner ist schneller da.", idx)
    notes(s, "Live-Demo (src/demo_paint.py auf dem DGX). Diese Folie baut in zwei Klicks auf: erst ein graues "
             "Feld mit '?' (Frage ans Publikum: was rechnet der Tuner da?), dann beim naechsten Klick das Bild. "
             "Idee: ein 1920x1080-Plasma-Bild ist die Summe von ~500 2D-Wellen, und diese Summe IST ein "
             "Matrixprodukt Y@X -> ein grosses, ganz normales GEMM (cmk,ckn->cmn), das der Tuner tunen kann. Wir "
             "rechnen es zweimal: mit der naiven 8x8-Default und mit der Config, die der Tuner waehlt (autotune, "
             "v2-Top-7). Beide gegen torch.einsum geprueft -> gleiches, korrektes Bild. Die Shape ist absichtlich "
             "krumm (1080/1920/1000 nicht tile-teilbar), damit die 8x8-Default schlecht padded und der Tuner "
             "sichtbar gewinnt (Padding 2048x2048 vs. ~1536x2048, typisch ~1.3-1.8x -- die echte Zahl zeigt der "
             "Live-Lauf). Das Skript zeigt live den Funnel (enumerate 486 -> prune -> Top-7), die zwei Zeiten als "
             "Balken, den Gewinner und speichert das Bild. Fallback ohne GPU: 'python demo_paint.py --preview' "
             "rechnet das Bild auf der CPU (das ist die Vorschau hier). Ehrlich: derselbe GEMM-Gewinn wie in der "
             "Eval, nur zum Anfassen. Wenn du LIVE ausfuehrst, das echte Bild/den echten Speedup statt der "
             "Vorschau zeigen.")


demo_slide(False)
demo_slide(True)

# --- 20 Fazit ---
idx += 1
s = slide(NAVY)
rect(s, 0, 0, 0.35, SH, fill=BLUE)
textbox(s, 0.9, 0.68, 11.5, 0.38, [[("FAZIT", 15, BLUE, True, False)]])
textbox(s, 0.86, 1.18, 12.0, 0.9, [[("Eingrenzen ist sicher — entscheiden muss man messen",
        26, WHITE, True, False)]])
rect(s, 0.9, 2.05, 2.0, 0.06, fill=ORANGE)
lite = RGBColor(0xd7, 0xe3, 0xf2)
bullets2 = [
    ("A05:", "Tuner reproduziert Hand (Ø 1.03×), gewinnt bei krummen Shapes (+58 %)."),
    ("A06:", "sauberer Transfer, ein zusätzlicher Kernel-Typ, schlägt den Handkernel (+29 %)."),
    ("cuBLAS:", "auf GEMM nicht geschlagen — Wert ist Allgemeinheit + Gewinne ohne guten Library-Pfad."),
    ("Modell:", "v2 bester Vorfilter (97.8 % @ top-7), roofline bester Ranker — messen bleibt entscheidend."),
    ("Praxis:", "GPU-spezifisch + gecacht — auf GB10 und 3070 bestätigt: Hebel wirkt, beste Config je GPU verschieden."),
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.6), Inches(4.4))
tf = tb.text_frame; tf.word_wrap = True
for i, (head, body) in enumerate(bullets2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = "▸  "; _set(r, 17, ORANGE, True)
    r1 = p.add_run(); r1.text = head + " "; _set(r1, 17, WHITE, True)
    r2 = p.add_run(); r2.text = body; _set(r2, 17, lite, False)
_reg(idx, "fazit-bullets", 0.9, 2.4, 11.6, 4.4)
notes(s, "Kernaussagen: A05 reproduziert Hand, A06 transferiert und schlaegt den Handkernel, cuBLAS wird auf "
         "GEMM nicht geschlagen (Wert = Allgemeinheit + Gewinne wo die Library keinen Pfad hat), das Modell "
         "taugt als Vorfilter (v2) aber nicht als exakter Ranker, und der Praxiswert ist GPU-spezifisches, "
         "gecachtes Tuning. Wichtigste Erkenntnis: die Eingrenzung ist das Sichere, die Entscheidung muss "
         "man messen.")

prs.save(OUT)
print("gespeichert:", OUT, "-", len(prs.slides._sldIdLst), "Folien")

# ---- Bounds-Check ----
issues = 0
for (i, name, L, T, W, H) in _placed:
    if L < -0.01 or T < -0.01 or L + W > SW + 0.01 or T + H > SH + 0.01:
        print(f"  !! ausserhalb: Folie {i+1} {name} L={L:.2f} T={T:.2f} R={L+W:.2f} B={T+H:.2f}")
        issues += 1
print(f"Bounds-Check: {len(_placed)} platzierte Objekte, {issues} ausserhalb der Folie.")
